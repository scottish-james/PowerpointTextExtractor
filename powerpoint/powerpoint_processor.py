"""
Main orchestrator for the PowerPoint-to-markdown pipeline.

Sits between the Flask application and the specialist components. When a file
comes in, this module decides how to process it, coordinates the other components
in the right order, and returns the finished markdown.

Pipeline (XML path):
    1. AccessibilityOrderExtractorV2 — reads shapes from each slide in reading order
    2. ContentExtractor             — pulls text, images, tables and charts from each shape
    3. TextProcessor                — cleans and formats raw text
    4. MarkdownConverter            — turns the structured data into markdown
    5. MetadataExtractor            — wraps the output with presentation metadata

If XML access is unavailable the file is passed directly to MarkItDown as a fallback.
"""

from pptx import Presentation
import os
from datetime import datetime
from markitdown import MarkItDown

from .accessibility_extractor_v2 import AccessibilityOrderExtractorV2
from .content_extractor import ContentExtractor
from .diagram_analyzer import DiagramAnalyzer
from .text_processor import TextProcessor
from .markdown_converter import MarkdownConverter
from .metadata_extractor import MetadataExtractor


class PowerPointProcessor:
    """
    Main PowerPoint processor with semantic role information flow and group handling.
    """

    def __init__(self, use_accessibility_order=True):
        self.use_accessibility_order = use_accessibility_order

        self.accessibility_extractor = AccessibilityOrderExtractorV2(use_accessibility_order)
        self.content_extractor = ContentExtractor()
        self.diagram_analyzer = DiagramAnalyzer()
        self.text_processor = TextProcessor()
        self.markdown_converter = MarkdownConverter()
        self.metadata_extractor = MetadataExtractor()

        self.markitdown = MarkItDown()
        self.supported_formats = ['.pptx', '.ppt']

    def _has_xml_access(self, file_path):
        """Check if XML-based processing is possible for the given file."""
        try:
            prs = Presentation(file_path)
            return len(prs.slides) > 0
        except Exception:
            return False

    def extract_slide_data(self, slide, slide_number):
        """Extract content from a single slide using the component pipeline."""
        shapes_with_roles = self.accessibility_extractor.get_slide_reading_order(slide, slide_number)

        slide_data = {
            "slide_number": slide_number,
            "content_blocks": [],
            "extraction_method": "semantic_accessibility_order_v2"
        }

        for shape, semantic_role in shapes_with_roles:
            block = self.content_extractor.extract_shape_content(
                shape,
                self.text_processor,
                self.accessibility_extractor,
                groups_already_expanded=True,
                semantic_role=semantic_role
            )
            if block:
                slide_data["content_blocks"].append(block)

        return slide_data

    def configure_extraction_method(self, use_accessibility_order):
        """Configure reading order extraction method."""
        self.use_accessibility_order = use_accessibility_order
        self.accessibility_extractor.accessibility_order = use_accessibility_order

    def convert_pptx_to_markdown_enhanced(self, file_path, convert_slide_titles=True):
        """
        Main entry point. Uses XML processing where available, falls back to MarkItDown.
        convert_slide_titles kept for compatibility — XML semantic roles control titles.
        """
        try:
            if self._has_xml_access(file_path):
                return self._sophisticated_xml_processing(file_path, convert_slide_titles)
            else:
                return self._simple_markitdown_processing(file_path)
        except Exception as e:
            raise Exception(f"Error processing PowerPoint file: {str(e)}")

    def _sophisticated_xml_processing(self, file_path, convert_slide_titles):
        """Full processing pipeline using XML extraction and semantic roles."""
        prs = Presentation(file_path)

        pptx_metadata = self.metadata_extractor.extract_pptx_metadata(prs, file_path)
        structured_data = self.extract_presentation_data(prs)

        markdown = self.markdown_converter.convert_structured_data_to_markdown(
            structured_data, convert_slide_titles=False
        )

        markdown_with_metadata = self.metadata_extractor.add_pptx_metadata(
            markdown, pptx_metadata
        )

        return markdown_with_metadata

    def _simple_markitdown_processing(self, file_path):
        """Fallback processing using MarkItDown."""
        try:
            result = self.markitdown.convert(file_path)

            try:
                markdown_content = result.markdown
            except AttributeError:
                try:
                    markdown_content = result.text_content
                except AttributeError:
                    raise Exception("Neither 'markdown' nor 'text_content' attribute found on result object")

            return f"\n<!-- Converted using MarkItDown fallback - XML not available -->\n{markdown_content}"

        except Exception as e:
            raise Exception(f"MarkItDown processing failed: {str(e)}")

    def extract_presentation_data(self, presentation):
        """Extract structured data from all slides."""
        data = {
            "total_slides": len(presentation.slides),
            "slides": []
        }

        for slide_idx, slide in enumerate(presentation.slides, 1):
            slide_data = self.extract_slide_data(slide, slide_idx)
            data["slides"].append(slide_data)

        return data

    def get_processing_summary(self, file_path):
        """Return a summary of how the file would be processed, without converting it."""
        try:
            has_xml = self._has_xml_access(file_path)

            summary = {
                "file_path": file_path,
                "has_xml_access": has_xml,
                "processing_method": "sophisticated_xml_with_semantic_roles_v2" if has_xml else "markitdown_fallback"
            }

            if has_xml:
                prs = Presentation(file_path)

                summary.update({
                    "slide_count": len(prs.slides),
                    "extraction_method": "accessibility_order_v2_with_semantic_roles",
                    "has_diagram_analysis": True,
                    "has_semantic_title_detection": True,
                    "slides_preview": []
                })

                for i, slide in enumerate(prs.slides[:3], 1):
                    shapes_with_roles = self.accessibility_extractor.get_slide_reading_order(slide, i)

                    title_count = sum(1 for _, role in shapes_with_roles if role == "title")
                    subtitle_count = sum(1 for _, role in shapes_with_roles if role == "subtitle")
                    content_count = sum(1 for _, role in shapes_with_roles if role == "content")

                    summary["slides_preview"].append({
                        "slide_number": i,
                        "shape_count": len(shapes_with_roles),
                        "title_shapes": title_count,
                        "subtitle_shapes": subtitle_count,
                        "content_shapes": content_count,
                        "has_text": any(
                            hasattr(shape, 'text_frame') and shape.text_frame
                            for shape, _ in shapes_with_roles
                        ),
                        "extraction_method": "semantic_accessibility_order_v2"
                    })
            else:
                summary.update({
                    "slide_count": "unknown",
                    "extraction_method": "markitdown_fallback",
                    "has_diagram_analysis": False,
                    "has_semantic_title_detection": False,
                    "note": "XML not available - using simple MarkItDown conversion"
                })

            return summary

        except Exception as e:
            return {"error": str(e)}


def convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles=True):
    """
    Convenience function for backward compatibility.
    convert_slide_titles kept for compatibility — XML semantic roles control titles.
    """
    processor = PowerPointProcessor()
    return processor.convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles)


def process_powerpoint_file(file_path, output_format="markdown", convert_slide_titles=True):
    """
    Convenience function for file processing with multiple output options.
    convert_slide_titles kept for compatibility — XML semantic roles control titles.
    """
    processor = PowerPointProcessor()

    if output_format == "summary":
        return processor.get_processing_summary(file_path)

    markdown_content = processor.convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles)

    result = {
        "content": markdown_content,
        "format": output_format,
        "processing_method": "sophisticated_xml_with_semantic_roles_v2" if processor._has_xml_access(
            file_path) else "markitdown_fallback"
    }

    if processor._has_xml_access(file_path):
        try:
            prs = Presentation(file_path)
            result["metadata"] = processor.metadata_extractor.extract_pptx_metadata(prs, file_path)
        except Exception:
            pass

    return result