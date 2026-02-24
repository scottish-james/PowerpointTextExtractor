"""
PowerPoint Reading Order Extractor V2 - Production Ready
Extracts shapes in PowerPoint slides in proper reading order with semantic roles.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET


class AccessibilityOrderExtractorV2:
    """
    Extracts shapes from PowerPoint slides in proper reading order with semantic roles.

    Args:
        use_accessibility_order: If True, uses semantic XML-based ordering. If False, uses basic recursive expansion.
    """

    def __init__(self, use_accessibility_order=True):
        self.accessibility_order = use_accessibility_order
        self.use_accessibility_order = use_accessibility_order  # Backward compatibility
        self.last_extraction_method = "not_extracted"

        # XML namespaces for PowerPoint OOXML processing
        self.namespaces = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        }

        # Store shape classifications for later retrieval
        self.shape_classifications = {}

    def get_slide_reading_order(self, slide, slide_number):
        """
        Main method to get reading order for a slide with proper role assignment.

        Returns:
            List of (shape, role) tuples in reading order
        """
        original_shapes = list(slide.shapes)

        if self.accessibility_order:
            final_shapes = self._get_semantic_accessibility_order(slide)
            self.last_extraction_method = "semantic_accessibility_order"
        else:
            final_shapes = self._expand_all_groups_recursively(original_shapes)
            self.last_extraction_method = "recursive_group_expansion"

        shapes_with_roles = []
        for shape in final_shapes:
            if hasattr(self, 'shape_classifications') and id(shape) in self.shape_classifications:
                semantic_role = self.shape_classifications[id(shape)]
            else:
                semantic_role = self._get_semantic_role_from_xml(shape)

            shapes_with_roles.append((shape, semantic_role))

        return shapes_with_roles

    def _get_semantic_accessibility_order(self, slide):
        """
        Semantic ordering with duplicate elimination.
        """
        xml_ordered_shapes = self._get_xml_document_order_deduplicated(slide)

        final_ordered_shapes = []
        for shape in xml_ordered_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_children = self._expand_all_groups_recursively([shape])
                final_ordered_shapes.extend(group_children)
            else:
                final_ordered_shapes.append(shape)

        deduplicated_shapes = self._deduplicate_shapes_by_object_id(final_ordered_shapes)

        title_shapes = []
        content_shapes = []
        other_shapes = []
        shape_classifications = {}

        for shape in deduplicated_shapes:
            shape_name = getattr(shape, 'name', '')

            text_preview = ""
            try:
                if hasattr(shape, 'text') and shape.text:
                    text_preview = shape.text.strip()
            except:
                pass

            if "title" in shape_name.lower() and "subtitle" not in shape_name.lower():
                title_shapes.append(shape)
                role = "title"
            elif "slide number" in shape_name.lower():
                role = "slide_number"
            elif text_preview or (hasattr(shape, 'has_table') and shape.has_table) or (hasattr(shape, 'has_chart') and shape.has_chart):
                content_shapes.append(shape)
                role = "content"
            else:
                other_shapes.append(shape)
                role = "other"

            shape_classifications[id(shape)] = role

        self.shape_classifications = shape_classifications

        return title_shapes + content_shapes + other_shapes

    def _expand_all_groups_recursively(self, shapes, depth=0):
        """
        Recursively expand all groups to extract individual shapes.
        """
        expanded_shapes = []

        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_children = list(shape.shapes)
                recursively_expanded = self._expand_all_groups_recursively(group_children, depth + 1)
                expanded_shapes.extend(recursively_expanded)
            else:
                expanded_shapes.append(shape)

        return expanded_shapes

    def _get_semantic_role_from_xml(self, shape):
        """
        Determine semantic role of a shape based on name and content.
        """
        try:
            shape_name = getattr(shape, 'name', '').lower()

            if "title" in shape_name and "subtitle" not in shape_name:
                return "title"
            elif "subtitle" in shape_name or "sub-title" in shape_name:
                return "subtitle"
            elif "slide number" in shape_name:
                return "slide_number"

            if hasattr(shape, 'text') and shape.text:
                text = shape.text.lower().strip()

                if len(text) < 100 and any(keyword in text for keyword in ['title', 'heading', 'header']):
                    return "title"
                elif any(keyword in text for keyword in ['subtitle', 'subheading', 'sub-title']):
                    return "subtitle"
                elif len(text) > 10:
                    return "content"
                else:
                    return "other"

            if hasattr(shape, 'shape_type'):
                if shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART]:
                    return "content"
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    return "content"

            return "other"

        except Exception:
            return "other"

    def _get_xml_document_order_deduplicated(self, slide):
        """
        Get shapes in actual XML document order with deduplication.
        """
        try:
            slide_element = slide._element
            ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
            shape_tree = slide_element.find('.//p:spTree', ns)

            if shape_tree is None:
                return self._deduplicate_shapes_by_object_id(list(slide.shapes))

            shape_elements = []
            for child in shape_tree:
                if child.tag.endswith('}sp') or child.tag.endswith('}grpSp') or child.tag.endswith(
                        '}pic') or child.tag.endswith('}cxnSp') or child.tag.endswith('}graphicFrame'):
                    shape_elements.append(child)

            def get_xml_id(elem):
                try:
                    for child in elem.iter():
                        if child.tag.endswith('}cNvPr'):
                            return child.get('id')
                except Exception:
                    pass
                return None

            id_to_shape_map = {}
            for shape in slide.shapes:
                xml_id = get_xml_id(shape._element)
                if xml_id is not None:
                    id_to_shape_map[xml_id] = shape

            ordered_shapes = []
            for xml_elem in shape_elements:
                xml_id = get_xml_id(xml_elem)
                if xml_id is not None and xml_id in id_to_shape_map:
                    ordered_shapes.append(id_to_shape_map[xml_id])

            return self._deduplicate_shapes_by_object_id(ordered_shapes)

        except Exception:
            return self._deduplicate_shapes_by_object_id(list(slide.shapes))

    def _deduplicate_shapes_by_object_id(self, shapes):
        """
        Remove duplicate shapes based on object ID.
        """
        seen_ids = set()
        deduplicated = []

        for shape in shapes:
            shape_id = id(shape)
            if shape_id not in seen_ids:
                seen_ids.add(shape_id)
                deduplicated.append(shape)

        return deduplicated

    def get_reading_order_of_grouped_shapes(self, group_shape):
        """
        Get reading order for shapes within a group.
        Required by ContentExtractor for processing group children.
        """
        try:
            return list(group_shape.shapes)
        except Exception:
            return []

    def get_reading_order_of_grouped_by_shape(self, shape):
        """
        Alias for backward compatibility.
        """
        return self.get_reading_order_of_grouped_shapes(shape)

    def get_last_extraction_method(self):
        """
        Return the last extraction method used.
        """
        return self.last_extraction_method

    def _has_xml_access(self, slide):
        """
        Check if XML access is available for the slide.
        """
        try:
            return len(slide.shapes) >= 0
        except:
            return False

    @property
    def use_accessibility_order(self):
        return self.accessibility_order

    @use_accessibility_order.setter
    def use_accessibility_order(self, value):
        self.accessibility_order = value