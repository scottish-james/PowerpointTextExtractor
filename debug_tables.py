"""
Drop this script next to a .pptx file with a table and run it.
It traces exactly what the pipeline sees at each stage.
"""

import sys
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SRC_DIR = os.path.join(SCRIPT_DIR, "src")
sys.path.insert(0, SRC_DIR)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

file_path = sys.argv[1]
prs = Presentation(file_path)

for slide_num, slide in enumerate(prs.slides, 1):
    print(f"\n=== Slide {slide_num} ===")
    print(f"  Total shapes from slide.shapes: {len(slide.shapes)}")

    for i, shape in enumerate(slide.shapes):
        shape_type = shape.shape_type
        shape_type_str = str(shape_type)
        shape_type_name = shape_type_str.split('.')[-1]
        has_table = hasattr(shape, 'has_table') and shape.has_table
        has_chart = hasattr(shape, 'has_chart') and shape.has_chart
        has_text = hasattr(shape, 'text') and bool(shape.text.strip()) if hasattr(shape, 'text') else False

        type_int = int(shape_type) if shape_type is not None else 'None'
        print(f"  Shape {i}: name='{shape.name}' | type_int={type_int} | type_str='{shape_type_str}' | type_name='{shape_type_name}' | has_table={has_table} | has_chart={has_chart} | has_text={has_text}")

        # Check what the tag is in the XML
        try:
            tag = shape._element.tag.split('}')[-1] if '}' in shape._element.tag else shape._element.tag
            print(f"           XML tag: {tag}")
        except:
            print(f"           XML tag: could not read")

    # Now check the spTree XML directly
    print(f"\n  --- spTree XML tags ---")
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        shape_tree = slide._element.find('.//p:spTree', ns)
        if shape_tree is not None:
            for child in shape_tree:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                sp_id = child.attrib.get('id', child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr') and child.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr').get('id', '?'))
                print(f"    tag={tag}")
        else:
            print("    spTree not found")
    except Exception as e:
        print(f"    Error reading spTree: {e}")